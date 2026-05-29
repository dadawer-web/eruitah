import io
import logging

import httpx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


def hex_to_rgb(hex_str: str) -> RGBColor:
    try:
        hex_str = str(hex_str).strip().lstrip('#')
        if len(hex_str) == 6:
            return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
    except Exception:
        pass
    return RGBColor(30, 30, 30)


async def generate_pptx_from_json(ppt_data: dict) -> io.BytesIO:
    prs = Presentation()

    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    design = ppt_data.get("design", {})
    palette = design.get("palette", {})
    bg_color = hex_to_rgb(palette.get("bg", palette.get("background", "#ffffff")))
    text_color = hex_to_rgb(palette.get("text", "#1f2937"))
    accent_color = hex_to_rgb(palette.get("accent", "#6366f1"))

    slides_data = ppt_data.get("slides", [])

    headers = {
        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
        "Accept": "image/*",
    }

    async with httpx.AsyncClient(verify=False) as client:
        for slide_data in slides_data:
            layout_type = slide_data.get("layout", "content")

            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)

            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = bg_color

            # ==================== Cover ====================
            if layout_type == "cover":
                title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(1.5))
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = str(slide_data.get("title") or "演示文稿")
                p.font.size = Pt(46)
                p.font.bold = True
                p.font.color.rgb = accent_color

                if "subtitle" in slide_data or "eyebrow" in slide_data:
                    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11.333), Inches(1.2))
                    stf = sub_box.text_frame
                    stf.word_wrap = True

                    if slide_data.get("eyebrow"):
                        ep = stf.paragraphs[0]
                        ep.text = str(slide_data.get("eyebrow") or "").upper()
                        ep.font.size = Pt(14)
                        ep.font.color.rgb = accent_color
                        ep.space_after = Pt(6)
                        sp = stf.add_paragraph()
                    else:
                        sp = stf.paragraphs[0]

                    sp.text = str(slide_data.get("subtitle") or "")
                    sp.font.size = Pt(22)
                    sp.font.color.rgb = text_color

                image_url = slide_data.get("image_url")
                if image_url and str(image_url).startswith("http"):
                    try:
                        resp = await client.get(image_url, headers=headers, timeout=20.0, follow_redirects=True)
                        if resp.status_code == 200:
                            img_stream = io.BytesIO(resp.content)
                            slide.shapes.add_picture(img_stream, Inches(4.5), Inches(4.8), width=Inches(4))
                    except Exception as e:
                        logger.warning(f"下载封面图片失败: {e}")

            # ==================== Section ====================
            elif layout_type == "section":
                box = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10.333), Inches(2.5))
                tf = box.text_frame
                tf.word_wrap = True

                p = tf.paragraphs[0]
                p.text = str(slide_data.get("title") or "")
                p.font.size = Pt(40)
                p.font.bold = True
                p.font.color.rgb = accent_color

                if slide_data.get("subtitle"):
                    sp = tf.add_paragraph()
                    sp.text = str(slide_data.get("subtitle") or "")
                    sp.font.size = Pt(20)
                    sp.font.color.rgb = text_color
                    sp.space_before = Pt(14)

                if slide_data.get("eyebrow"):
                    eb = tf.add_paragraph()
                    eb.text = str(slide_data.get("eyebrow") or "").upper()
                    eb.font.size = Pt(12)
                    eb.font.color.rgb = accent_color
                    eb.space_before = Pt(8)

            # ==================== Closing ====================
            elif layout_type == "closing":
                box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.333), Inches(3.0))
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = str(slide_data.get("title") or "Thank You")
                p.font.size = Pt(48)
                p.font.bold = True
                p.font.color.rgb = accent_color
                p.alignment = 1  # PP_ALIGN.CENTER

            # ==================== Content / Code ====================
            else:
                title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.333), Inches(1.0))
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = str(slide_data.get("title") or "")
                p.font.size = Pt(32)
                p.font.bold = True
                p.font.color.rgb = accent_color

                image_url = slide_data.get("image_url")
                has_valid_img = image_url and str(image_url).startswith("http") and layout_type != "closing"

                if has_valid_img:
                    text_left, text_width = Inches(1.0), Inches(7.2)
                    img_left, img_width = Inches(8.6), Inches(3.8)
                else:
                    text_left, text_width = Inches(1.0), Inches(11.333)

                content_box = slide.shapes.add_textbox(text_left, Inches(1.8), text_width, Inches(5.0))
                ctf = content_box.text_frame
                ctf.word_wrap = True

                first_para = True
                for comp in slide_data.get("components") or []:
                    comp_type = comp.get("type")
                    content = str(comp.get("content") or "").strip()

                    if comp_type in ["text", "heading", "card"]:
                        p = ctf.add_paragraph() if not first_para else ctf.paragraphs[0]
                        p.text = content
                        p.font.color.rgb = text_color
                        p.space_after = Pt(6)
                        if comp_type == "heading":
                            p.font.bold = True
                            p.font.size = Pt(24)
                            p.space_after = Pt(10)
                        elif comp_type == "card":
                            p.font.size = Pt(18)
                            p.font.italic = True
                        else:
                            p.font.size = Pt(18)
                        first_para = False

                    elif comp_type == "bullet_list":
                        items = comp.get("items") or []
                        for item in items:
                            p = ctf.add_paragraph() if not first_para else ctf.paragraphs[0]
                            p.text = str(item or "")
                            p.font.size = Pt(18)
                            p.font.color.rgb = text_color
                            p.level = 1
                            p.space_after = Pt(4)
                            first_para = False

                    elif comp_type == "code":
                        p = ctf.add_paragraph() if not first_para else ctf.paragraphs[0]
                        lang = str(comp.get("language") or "code")
                        p.text = f"[{lang}]\n{content}"
                        p.font.name = "Courier New"
                        p.font.color.rgb = RGBColor(120, 120, 120)
                        p.font.size = Pt(14)
                        p.space_after = Pt(8)
                        first_para = False

                    elif comp_type == "divider":
                        p = ctf.add_paragraph() if not first_para else ctf.paragraphs[0]
                        p.text = "─────────────────────────────"
                        p.font.size = Pt(10)
                        p.font.color.rgb = accent_color
                        p.space_after = Pt(8)
                        first_para = False

                    elif comp_type == "image":
                        comp_img_url = comp.get("image_url")
                        if comp_img_url and str(comp_img_url).startswith("http"):
                            try:
                                resp = await client.get(comp_img_url, headers=headers, timeout=15.0, follow_redirects=True)
                                if resp.status_code == 200:
                                    img_stream = io.BytesIO(resp.content)
                                    slide.shapes.add_picture(img_stream, Inches(1.0), Inches(5.5), width=Inches(3))
                            except Exception:
                                pass

                if has_valid_img:
                    try:
                        resp = await client.get(image_url, headers=headers, timeout=20.0, follow_redirects=True)
                        if resp.status_code == 200:
                            img_stream = io.BytesIO(resp.content)
                            slide.shapes.add_picture(img_stream, img_left, Inches(2.0), width=img_width)
                    except Exception as e:
                        logger.warning(f"下载内容页图片失败: {e}")

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream
