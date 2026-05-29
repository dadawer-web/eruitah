import asyncio
import base64
import io
import json
import logging
import re
from typing import Any, Dict, List, Optional

from app.models.ppt_schema import (
    DesignPalette,
    DesignSystem,
    DesignTypeScale,
    Presentation,
    SlideComponent,
    SlideMeta,
    SlidePage,
)

logger = logging.getLogger(__name__)

PPT_GENERATION_SYSTEM_PROMPT = """你是一个专业的演示文稿设计专家。你需要根据用户的需求，生成一份完整的 PPT 结构化 JSON。

## 输出格式

请严格输出以下 JSON 结构（不要包含任何 markdown 标记）：

{
  "meta": {
    "title": "演示文稿标题",
    "theme": "主题名称"
  },
  "design": {
    "palette": {"bg": "#f7f5f0", "text": "#1a1814", "accent": "#6d4cff"},
    "fonts": {
      "display": "Georgia, serif",
      "body": "-apple-system, BlinkMacSystemFont, system-ui, sans-serif"
    },
    "type_scale": {"hero": 168, "body": 36},
    "radius": 12
  },
  "slides": [
    {
      "layout": "cover",
      "title": "封面标题",
      "subtitle": "副标题",
      "eyebrow": "标签文字",
      "image_prompt": "modern business corporate background abstract",
      "image_search_keyword": "modern glass tower skyline blue",
      "image_visual_description": "现代玻璃幕墙摩天大楼在蓝色天空下的倒影，商务感十足",
      "components": []
    },
    {
      "layout": "content",
      "title": "内容页标题",
      "image_prompt": "professional teamwork office meeting",
      "image_search_keyword": "people high five office bright",
      "image_visual_description": "办公室里多人击掌庆祝，明亮自然光，团队协作氛围",
      "components": [
        {"type": "text", "content": "正文内容"},
        {"type": "bullet_list", "items": ["要点1", "要点2", "要点3"]},
        {"type": "code", "content": "print('hello')", "language": "python"}
      ]
    }
  ]
}

## 可用的 layout 类型

- **cover**: 封面页，包含大标题、副标题、eyebrow
- **section**: 章节分隔页，大标题居中
- **content**: 标准内容页，标题 + 组件列表
- **two_column**: 双栏布局
- **code_focus**: 代码展示页，突出代码块
- **image_text**: 图文混排页
- **closing**: 结尾页

## 可用的 component 类型

- **heading**: 标题文本，content 为标题内容
- **text**: 普通文本，content 为文本内容
- **code**: 代码块，content 为代码，language 为语言名
- **bullet_list**: 列表，items 为列表项数组
- **image**: 图片，image_url 为图片地址
- **divider**: 分隔线
- **card**: 卡片，content 为卡片内容，accent 为强调色
- **two_column**: 双栏，content 为左栏，items 为右栏内容列表

## 设计原则

1. 封面页必须有 eyebrow 和 subtitle
2. 每页组件数量控制在 2-5 个
3. 代码页使用 code_focus 布局
4. 结尾页使用 closing 布局
5. 颜色使用 hex 格式
6. 内容要充实、专业，不要敷衍"""


PPT_GENERATION_USER_TEMPLATE = """请根据以下需求生成一份 PPT：

## 用户需求
{user_instruction}

## PPT 分析数据（如有）
{ppt_analysis}

请生成完整的 JSON 结构，确保幻灯片数量合理（通常 5-15 页），内容专业详实。"""


class PPTGenerator:
    def __init__(self, llm_client=None):
        self.llm_client = llm_client
        self.analyzer = PPTAnalyzer()

    async def generate_slide_deck(
        self,
        user_instruction: str,
        ppt_file_path: Optional[str] = None,
        rag_context: str = "",
    ) -> Dict[str, Any]:
        ppt_analysis = ""
        if ppt_file_path:
            try:
                analysis = await self.analyzer.extract_text_and_structure(ppt_file_path)
                ppt_analysis = json.dumps(analysis, ensure_ascii=False, indent=2)
            except Exception as e:
                logger.warning(f"PPT analysis failed: {e}")
                ppt_analysis = "（无法解析 PPT 文件）"

        if not self.llm_client:
            return self._build_fallback_deck(user_instruction)

        rag_section = ""
        if rag_context:
            rag_section = f"\n\n【知识库参考】（生成PPT时必须融合以下专业知识和公司信息）:\n{rag_context}"

        messages = [
            {"role": "system", "content": PPT_GENERATION_SYSTEM_PROMPT},
            {
                "role": "user",
                "content": PPT_GENERATION_USER_TEMPLATE.format(
                    user_instruction=user_instruction,
                    ppt_analysis=ppt_analysis or "（无 PPT 文件）",
                ) + rag_section,
            },
        ]

        try:
            logger.info("🚀 [PPTAgent] 正在调用大模型生成 SlideDeck JSON...")
            response = await self.llm_client.acall_api(messages)
            if not response:
                logger.error("🚨 [PPTAgent] 大模型返回空响应，使用兜底模板")
                return self._build_fallback_deck(user_instruction)

            logger.info(f"📋 [PPTAgent] 大模型返回内容长度: {len(response)}")
            deck = self._parse_deck_response(response)
            if deck:
                return deck.model_dump()
            logger.warning("🚨 [PPTAgent] JSON 解析失败，使用兜底模板")
            return self._build_fallback_deck(user_instruction)
        except Exception as e:
            logger.error(f"🚨 [PPTAgent] 生成失败: {repr(e)}")
            return self._build_fallback_deck(user_instruction)

    def _parse_deck_response(self, response: str) -> Optional[Presentation]:
        cleaned = response.strip()

        code_block = re.search(r"```(?:json)?\s*\n?(.*?)\n?\s*```", cleaned, re.DOTALL)
        if code_block:
            cleaned = code_block.group(1).strip()

        brace_match = re.search(r"\{.*\}", cleaned, re.DOTALL)
        if brace_match:
            cleaned = brace_match.group(0)

        try:
            parsed = json.loads(cleaned)
            return Presentation(**parsed)
        except Exception:
            pass

        try:
            fixed = re.sub(r"'", '"', cleaned)
            fixed = re.sub(r',\s*}', '}', fixed)
            fixed = re.sub(r',\s*]', ']', fixed)
            parsed = json.loads(fixed)
            return Presentation(**parsed)
        except Exception:
            pass

        return None

    def _build_fallback_deck(self, instruction: str) -> Dict[str, Any]:
        deck = Presentation(
            meta=SlideMeta(title=instruction[:50] if instruction else "Untitled"),
            design=DesignSystem(),
            slides=[
                SlidePage(
                    layout="cover",
                    title=instruction[:60] if instruction else "AI Generated Presentation",
                    subtitle="Powered by Document Copilot",
                    eyebrow="AI 自动生成",
                    image_prompt="futuristic technology abstract blue glow",
                    image_search_keyword="futuristic robot hand network",
                    image_visual_description="发光的机械手触碰全息网络图，深蓝色背景，科技感十足",
                ),
                SlidePage(
                    layout="content",
                    title="概述",
                    components=[
                        SlideComponent(
                            type="text",
                            content=f"本演示文稿根据以下需求自动生成：{instruction}",
                        ),
                    ],
                    image_prompt="document paper stack desk",
                    image_search_keyword="document paper stack desk",
                    image_visual_description="整齐叠放的文件和纸张在木质桌面上，柔和自然光",
                ),
                SlidePage(
                    layout="closing",
                    title="谢谢",
                    components=[],
                ),
            ],
        )
        return deck.model_dump()


class PPTAnalyzer:
    async def extract_text_and_structure(
        self, file_path: str
    ) -> Dict[str, Any]:
        result = await asyncio.to_thread(self._extract_sync, file_path)
        logger.info(
            f"PPT analyzed: {file_path}, slides={result['slide_count']}"
        )
        return result

    async def extract_slide(self, file_path: str, slide_index: int) -> Dict[str, Any]:
        result = await self.extract_text_and_structure(file_path)
        slides = result.get("slides", [])
        if 0 <= slide_index < len(slides):
            return slides[slide_index]
        return {"error": f"Slide index {slide_index} out of range"}

    async def get_outline(self, file_path: str) -> List[Dict[str, Any]]:
        result = await self.extract_text_and_structure(file_path)
        return [
            {"index": s["index"], "title": s.get("title", ""), "layout": s.get("layout_name", "")}
            for s in result.get("slides", [])
        ]

    def _extract_sync(self, file_path: str) -> Dict[str, Any]:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation(file_path)
        slides_data = []
        all_images = []
        image_idx = 0

        for idx, slide in enumerate(prs.slides):
            title = ""
            texts = []
            has_table = False
            slide_image_placeholders = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_text = shape.text_frame.text.strip()
                    if shape_text:
                        texts.append(shape_text)
                    if not title and shape.shape_type == 14:
                        title = shape_text

                if shape.has_table:
                    has_table = True

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        blob = image.blob
                        ct = image.content_type
                        if not ct:
                            ct = "image/png"
                        b64 = base64.b64encode(blob).decode("utf-8")
                        data_uri = f"data:{ct};base64,{b64}"
                        all_images.append(data_uri)
                        placeholder = f"[Slide_{idx + 1}_Image_{image_idx}]"
                        slide_image_placeholders.append(placeholder)
                        image_idx += 1
                    except Exception as e:
                        logger.warning(f"PPT: failed to extract image from slide {idx + 1}: {e}")

            if not title and texts:
                title = texts[0][:100]

            layout_name = ""
            try:
                layout_name = slide.slide_layout.name if slide.slide_layout else ""
            except Exception:
                pass

            if slide_image_placeholders:
                texts.append("文档图片: " + ", ".join(slide_image_placeholders))

            slides_data.append(
                {
                    "index": idx,
                    "title": title,
                    "texts": texts,
                    "has_table": has_table,
                    "has_image": len(slide_image_placeholders) > 0,
                    "layout_name": layout_name,
                }
            )

        full_text = "\n\n".join(
            f"[Slide {s['index'] + 1}] {s['title']}\n" + "\n".join(s["texts"])
            for s in slides_data
        )

        return {
            "slide_count": len(slides_data),
            "slides": slides_data,
            "full_text": full_text,
            "images": all_images,
        }
