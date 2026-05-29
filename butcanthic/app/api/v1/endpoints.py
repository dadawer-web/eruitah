import asyncio
import json
import logging
import os
import time
import uuid
from datetime import datetime
from typing import Dict, List, Optional

from fastapi import APIRouter, BackgroundTasks, Body, Depends, File, Form, Header, HTTPException, Query, UploadFile, status
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.security import OAuth2PasswordRequestForm
from pydantic import BaseModel, Field
from sqlalchemy.orm import Session

import redis as redis_lib

from app.core.security import (
    create_access_token,
    get_current_user,
    get_current_user_optional,
    get_current_user_from_header_or_query,
    verify_password,
    get_password_hash,
)

from app.models.ppt_schema import SlidePage

from app.models.database import DocumentMeta, User, init_db, get_session, get_db

from app.models.schemas import (
    DocumentProcessRequest,
    DocumentProcessResponse,
    DocumentUploadResponse,
    FlashcardResponse,
    HealthResponse,
    KnowledgeRecord,
    KnowledgeUploadRequest,
    KnowledgeUploadResponse,
    ReviewSubmit,
    SupportedFileType,
    TaskProgress,
    TaskStatus,
)

logger = logging.getLogger(__name__)


def clean_and_parse_json(raw_text: str) -> dict:
    import re as _re

    text = _re.sub(r'<think[^>]*>.*?</think\s*>', '', raw_text, flags=_re.DOTALL)
    text = _re.sub(r'```(?:json)?\s*\n?', '', text)
    text = _re.sub(r'\n?\s*```', '', text)

    start = text.find('{')
    end = text.rfind('}')
    if start == -1 or end == -1 or end <= start:
        logger.error(f"clean_and_parse_json: no JSON object found, raw (first 500): {raw_text[:500]}")
        raise ValueError("No JSON object found in LLM response")

    json_str = text[start:end + 1]

    try:
        return json.loads(json_str)
    except json.JSONDecodeError:
        pass

    fixed = _re.sub(r"'", '"', json_str)
    fixed = _re.sub(r',\s*}', '}', fixed)
    fixed = _re.sub(r',\s*]', ']', fixed)
    try:
        return json.loads(fixed)
    except json.JSONDecodeError as e:
        logger.error(f"clean_and_parse_json: JSONDecodeError after fix: {e}, raw (first 500): {raw_text[:500]}")
        raise


router = APIRouter()

_tasks: Dict[str, Dict] = {}

_start_time = time.time()


class EditSlideRequest(BaseModel):
    slide_index: int
    instruction: str
    original_slide: dict


class UserRegisterRequest(BaseModel):
    username: str = Field(..., min_length=2, max_length=50)
    password: str = Field(..., min_length=6, max_length=72)


class SilentLoginRequest(BaseModel):
    user_id: int = Field(..., gt=0)


def _get_document_service():
    from app.core.app_state import app_state
    return app_state.document_service


def _get_rag_engine():
    from app.core.app_state import app_state
    return app_state.rag_engine


def _get_llm_client():
    from app.core.app_state import app_state
    return app_state.llm_client


def _bg_extract_memory(user_instruction: str, user_id: str, llm_client, rag_engine):
    try:
        import asyncio
        from app.services.memory_service import memory_manager
        loop = asyncio.new_event_loop()
        try:
            loop.run_until_complete(
                memory_manager.extract_and_store_memory(
                    user_instruction=user_instruction,
                    user_id=user_id,
                    llm_client=llm_client,
                    rag_engine=rag_engine,
                )
            )
        finally:
            loop.close()
    except Exception as e:
        logger.warning(f"Background memory extraction failed: {e}")


@router.get("/health", response_model=HealthResponse)
async def health_check():
    return HealthResponse(version="1.0.0", uptime=time.time() - _start_time)


@router.post("/auth/silent-login")
async def silent_login(req: SilentLoginRequest, db: Session = Depends(get_db)):
    qt_username = f"qt_user_{req.user_id}"
    user = db.query(User).filter(User.username == qt_username).first()
    if not user:
        user = User(
            id=str(uuid.uuid4()),
            username=qt_username,
            hashed_password=get_password_hash(f"qt_{req.user_id}_{uuid.uuid4().hex[:8]}"),
        )
        db.add(user)
        db.commit()
        db.refresh(user)
        logger.info(f"[SilentAuth] Auto-created user: {qt_username} (id={user.id})")
    token = create_access_token(user_id=user.id)
    return {"access_token": token, "token_type": "bearer", "user_id": user.id, "username": user.username}


@router.post("/auth/register")
async def register(user_in: UserRegisterRequest, db: Session = Depends(get_db)):
    if db.query(User).filter(User.username == user_in.username).first():
        raise HTTPException(status_code=400, detail="用户名已被注册")
    if len(user_in.username) < 2 or len(user_in.username) > 32:
        raise HTTPException(status_code=400, detail="用户名长度需在 2-32 之间")
    if len(user_in.password) < 6:
        raise HTTPException(status_code=400, detail="密码长度不能少于 6 位")
    new_user = User(
        id=str(uuid.uuid4()),
        username=user_in.username,
        hashed_password=get_password_hash(user_in.password),
    )
    db.add(new_user)
    db.commit()
    db.refresh(new_user)
    token = create_access_token(user_id=new_user.id)
    return {"access_token": token, "token_type": "bearer", "user_id": new_user.id, "username": new_user.username}


@router.post("/auth/login")
async def login(form_data: OAuth2PasswordRequestForm = Depends(), db: Session = Depends(get_db)):
    user = db.query(User).filter(User.username == form_data.username).first()
    if not user or not verify_password(form_data.password, user.hashed_password):
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="用户名或密码错误",
            headers={"WWW-Authenticate": "Bearer"},
        )
    token = create_access_token(user_id=user.id)
    return {"access_token": token, "token_type": "bearer", "user_id": user.id, "username": user.username}


@router.post("/auth/token")
async def create_token(user_id: str = Form(...)):
    token = create_access_token(user_id=user_id)
    return {"access_token": token, "token_type": "bearer", "user_id": user_id}


@router.get("/auth/me")
async def get_me(current_user_id: str = Depends(get_current_user), db: Session = Depends(get_db)):
    user = db.query(User).filter(User.id == current_user_id).first()
    if not user:
        raise HTTPException(status_code=404, detail="用户不存在")
    return {"user_id": user.id, "username": user.username, "created_at": user.created_at.isoformat() if user.created_at else None}


@router.get("/documents/list")
async def list_documents(
    page: int = 1,
    page_size: int = 20,
    status: Optional[str] = None,
    file_type: Optional[str] = None,
    current_user_id: str = Depends(get_current_user),
):
    session = get_session()
    try:
        query = session.query(DocumentMeta).filter(DocumentMeta.user_id == current_user_id)
        if status:
            query = query.filter(DocumentMeta.status == status)
        if file_type:
            query = query.filter(DocumentMeta.file_type == file_type)
        total = query.count()
        items = query.order_by(DocumentMeta.upload_time.desc()).offset((page - 1) * page_size).limit(page_size).all()
        return {
            "total": total,
            "page": page,
            "page_size": page_size,
            "items": [item.to_dict() for item in items],
        }
    finally:
        session.close()


@router.post("/document/upload", response_model=DocumentUploadResponse)
async def upload_document(
    file: UploadFile = File(...),
    current_user_id: Optional[str] = Depends(get_current_user_optional),
    x_user_id: Optional[str] = Header(None, alias="X-User-Id"),
):
    file_type = SupportedFileType.from_filename(file.filename or "")
    if file_type is None:
        raise HTTPException(status_code=400, detail=f"不支持的文件格式，仅支持 docx/xlsx/pptx")

    file_content = await file.read()
    file_size = len(file_content)
    if file_size > 50 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="文件大小超过50MB限制")

    from app.core.file_manager import file_manager
    effective_user_id = x_user_id or current_user_id or None
    file_path = file_manager.save_upload(file.filename, file_content, file_type.value, user_id=effective_user_id)

    task_id = uuid.uuid4().hex[:12]
    _tasks[task_id] = {
        "task_id": task_id,
        "status": TaskStatus.PENDING,
        "progress": 0.0,
        "stage": "文件已上传",
        "filename": file.filename,
        "file_type": file_type,
        "file_path": file_path,
        "file_size": file_size,
        "tables_total": 0,
        "tables_completed": 0,
        "error_message": None,
        "output_path": None,
        "created_at": datetime.now(),
        "updated_at": datetime.now(),
        "user_id": current_user_id or "",
    }

    logger.info(f"Upload: task={task_id}, type={file_type.value}, file={file.filename}, user={current_user_id or 'anonymous'}")

    return DocumentUploadResponse(
        task_id=task_id,
        filename=file.filename,
        file_type=file_type,
        file_size=file_size,
    )


@router.post("/document/process", response_model=DocumentProcessResponse)
async def process_document(
    request: DocumentProcessRequest,
    current_user_id: Optional[str] = Depends(get_current_user_optional),
    x_user_id: Optional[str] = Header(None, alias="X-User-Id"),
):
    task_id = request.task_id
    if task_id not in _tasks:
        raise HTTPException(status_code=404, detail=f"任务 {task_id} 不存在")

    task = _tasks[task_id]
    if task["status"] not in (TaskStatus.PENDING,):
        raise HTTPException(status_code=409, detail=f"任务状态为 {task['status']}，无法重复处理")

    effective_user_id = x_user_id or current_user_id or None

    doc_service = _get_document_service()
    if doc_service is None:
        raise HTTPException(status_code=503, detail="文档处理服务未初始化")

    task["status"] = TaskStatus.PROCESSING
    task["updated_at"] = datetime.now()

    async def progress_callback(stage: str, progress: float):
        task["stage"] = stage
        task["progress"] = progress
        task["updated_at"] = datetime.now()

    async def _run():
        try:
            result = await doc_service.process_document(
                task_id=task_id,
                input_path=task["file_path"],
                selected_tables=request.selected_tables,
                max_retries=request.max_retries,
                progress_callback=progress_callback,
            )
            task["status"] = TaskStatus.COMPLETED
            task["progress"] = 100.0
            task["output_path"] = result["output_path"]
            task["tables_total"] = len(result["tables"])
            task["tables_completed"] = sum(1 for t in result["tables"] if not t.get("error_message"))
            task["updated_at"] = datetime.now()
        except Exception as e:
            logger.error(f"[{task_id}] Failed: {e}", exc_info=True)
            task["status"] = TaskStatus.FAILED
            task["error_message"] = str(e)
            task["updated_at"] = datetime.now()

    asyncio.create_task(_run())

    return DocumentProcessResponse(task_id=task_id, status=TaskStatus.PROCESSING)


@router.post("/task/stream-process")
async def stream_process(
    background_tasks: BackgroundTasks,
    files: Optional[List[UploadFile]] = File(None),
    user_instruction: str = Form(""),
    model: Optional[str] = Form(None),
    max_retries: int = Form(3),
    thread_id: Optional[str] = Form(None),
    current_user_id: Optional[str] = Depends(get_current_user_optional),
    x_user_id: Optional[str] = Header(None, alias="X-User-Id"),
):
    """
    SSE 流式处理接口
    上传文件后实时推送 LangGraph 工作流的每一步状态
    支持多文件同时上传与跨文档协同分析
    支持无文件纯文本模式（仅提供 user_instruction 即可生成 PPT）
    支持 thread_id 会话保持，实现跨轮对话记忆
    登录用户自动将 user_id 注入工作流状态，确保 RAG 检索在专属 Collection 中执行
    """
    import uuid

    effective_thread_id = thread_id or f"thread_{uuid.uuid4().hex[:8]}"

    uploaded_files_info = []
    if files:
        from app.core.file_manager import file_manager
        effective_user_id = x_user_id or current_user_id or None
        for f in files:
            if not f.filename:
                continue
            file_type_val = SupportedFileType.from_filename(f.filename)
            if file_type_val is None:
                continue
            file_content = await f.read()
            if len(file_content) > 200 * 1024 * 1024:
                raise HTTPException(status_code=413, detail=f"文件 {f.filename} 大小超过200MB限制")
            path = file_manager.save_upload(f.filename, file_content, file_type_val.value, user_id=effective_user_id)
            uploaded_files_info.append({
                "path": path,
                "type": file_type_val.value,
                "filename": f.filename,
            })

    rag_engine = _get_rag_engine()
    llm_client = _get_llm_client()

    if rag_engine is None or llm_client is None:
        raise HTTPException(status_code=503, detail="AI服务未初始化，请检查配置")

    if current_user_id and user_instruction:
        try:
            from app.services.memory_service import memory_manager
            background_tasks.add_task(
                _bg_extract_memory,
                user_instruction=user_instruction,
                user_id=current_user_id,
                llm_client=llm_client,
                rag_engine=rag_engine,
            )
        except Exception as e:
            logger.warning(f"Memory background task scheduling failed: {e}")

    async def event_generator():
        from app.agent_workflow.graph import run_workflow_streaming

        yield f'data: {{"status": "session_info", "thread_id": "{effective_thread_id}"}}\n\n'

        async for event in run_workflow_streaming(
            uploaded_files=uploaded_files_info,
            rag_engine=rag_engine,
            llm_client=llm_client,
            user_instruction=user_instruction,
            max_retries=max_retries,
            thread_id=effective_thread_id,
            user_id=current_user_id or "",
        ):
            if event.get("status") == "ppt_ready":
                ppt_data_dict = event.get("ppt_data", {})
                ppt_data_json_str = json.dumps(ppt_data_dict, ensure_ascii=False, default=str)
                yield f'data: {{"status": "ppt_ready", "ppt_data": {ppt_data_json_str}}}\n\n'
            else:
                sse_data = json.dumps(event, ensure_ascii=False, default=str)
                yield f"data: {sse_data}\n\n"

            if event.get("status") in ("error", "success"):
                break

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


@router.get("/documents/my")
async def my_documents(
    page: int = 1,
    page_size: int = 20,
    current_user_id: str = Depends(get_current_user),
    db: Session = Depends(get_db),
):
    offset = (page - 1) * page_size
    query = db.query(DocumentMeta).filter(DocumentMeta.user_id == current_user_id)
    total = query.count()
    docs = query.order_by(DocumentMeta.upload_time.desc()).offset(offset).limit(page_size).all()
    return {
        "total": total,
        "page": page,
        "page_size": page_size,
        "documents": [doc.to_dict() for doc in docs],
    }


@router.post("/task/edit-slide")
async def edit_slide_endpoint(req: EditSlideRequest):
    llm_client = _get_llm_client()
    if llm_client is None or llm_client.langchain_llm is None:
        raise HTTPException(status_code=503, detail="AI服务未初始化")

    from langchain_core.prompts import ChatPromptTemplate

    llm = llm_client.langchain_llm
    structured_llm = llm.with_structured_output(SlidePage)

    prompt = ChatPromptTemplate.from_messages([
        ("system", "你是一个顶级的幻灯片编辑专家。你的任务是根据用户的修改指令，在保持原有核心结构和配图要求的前提下，只修改指定的幻灯片单页 JSON。必须严格返回修改后的 SlidePage 对象，不允许捏造非法的组件类型（只能用 heading, text, image, bullet_list, code, divider, card, two_column）。image_prompt 字段必须保留且不能为空。"),
        ("user", "【原幻灯片数据】:\n{original_slide}\n\n【用户修改指令】:\n{instruction}\n\n请输出修改后的单页 JSON：")
    ])

    chain = prompt | structured_llm
    try:
        new_slide = await chain.ainvoke({
            "original_slide": json.dumps(req.original_slide, ensure_ascii=False, indent=2),
            "instruction": req.instruction,
        })
        return {"status": "success", "slide_index": req.slide_index, "new_slide": new_slide.model_dump()}
    except Exception as e:
        logger.error(f"edit_slide failed: {e}", exc_info=True)
        return {"status": "error", "message": str(e)}


@router.post("/task/export-pptx")
async def export_pptx_endpoint(req: dict):
    from app.services.ppt_export_service import generate_pptx_from_json

    try:
        ppt_stream = await generate_pptx_from_json(req)
        return StreamingResponse(
            ppt_stream,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=Generated_Presentation.pptx"},
        )
    except Exception as e:
        logger.error(f"Failed to export PPTX: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/document/download/{task_id}")
async def download_document(task_id: str):
    if task_id not in _tasks:
        raise HTTPException(status_code=404, detail=f"任务 {task_id} 不存在")

    task = _tasks[task_id]
    if task["status"] != TaskStatus.COMPLETED:
        raise HTTPException(status_code=409, detail=f"任务状态为 {task['status']}")

    output_path = task.get("output_path")
    if not output_path or not os.path.exists(output_path):
        raise HTTPException(status_code=404, detail="输出文件不存在")

    return FileResponse(
        path=output_path,
        filename=os.path.basename(output_path),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


@router.get("/files/output/{filename}")
async def download_output_file(filename: str, x_user_id: Optional[str] = Header(None, alias="X-User-Id")):
    from app.core.file_manager import file_manager
    file_path = file_manager.get_output_path(filename, user_id=x_user_id)
    if not file_path:
        raise HTTPException(status_code=404, detail="文件不存在")
    return FileResponse(path=file_path, filename=filename)


@router.post("/knowledge/upload", response_model=KnowledgeUploadResponse)
async def upload_knowledge(
    request: KnowledgeUploadRequest,
    records: list[KnowledgeRecord],
    current_user_id: str = Depends(get_current_user),
):
    rag_engine = _get_rag_engine()
    if rag_engine is None:
        raise HTTPException(status_code=503, detail="RAG引擎未初始化")

    data = {"records": [r.data for r in records]}
    result = await rag_engine.ingest_data(data, request.collection_name, user_id=current_user_id)

    return KnowledgeUploadResponse(
        collection_name=result["collection"],
        record_count=result["ingested_count"],
    )


@router.post("/knowledge/upload-file")
async def upload_knowledge_file(
    file: UploadFile = File(...),
    current_user_id: str = Depends(get_current_user),
):
    """
    上传知识库文件，支持 .txt, .jsonl, .pdf, .docx, .md
    支持智能文本分块（Text Splitting），避免文档过长导致超出大模型上下文
    数据存入用户专属集合 (kb_{user_id})，实现 Collection 物理隔离
    """
    filename = file.filename or ""
    supported_exts = (".txt", ".jsonl", ".pdf", ".docx", ".md")
    if not filename.lower().endswith(supported_exts):
        raise HTTPException(status_code=400, detail=f"不支持的格式，当前支持: {', '.join(supported_exts)}")

    rag_engine = _get_rag_engine()
    if rag_engine is None:
        raise HTTPException(status_code=503, detail="RAG引擎未初始化")

    content = await file.read()

    if filename.lower().endswith(".jsonl"):
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("gbk", errors="ignore")

        records = []
        for line in text.strip().split("\n"):
            line = line.strip()
            if line:
                try:
                    records.append(json.loads(line))
                except json.JSONDecodeError:
                    continue

        if not records:
            raise HTTPException(status_code=400, detail="JSONL 文件中无有效记录")

        result = await rag_engine.ingest_data({"records": records}, "default", user_id=current_user_id)
        logger.info(f"Knowledge uploaded: {filename} → {result['ingested_count']} records in [{result['collection']}]")
        return {
            "status": "success",
            "collection_name": result["collection"],
            "record_count": result["ingested_count"],
            "filename": filename,
        }

    import tempfile
    from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
    from langchain_core.documents import Document
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    ext = os.path.splitext(filename)[1].lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(content)
        tmp_path = tmp.name

    try:
        documents = []
        if ext == ".pdf":
            loader = PyPDFLoader(tmp_path)
            documents = loader.load()
            try:
                from app.services.ocr_helper import ocr_pdf_images_v2
                ocr_text = await ocr_pdf_images_v2(tmp_path, ai_client=_get_llm_client())
                if ocr_text:
                    documents.append(Document(page_content=ocr_text, metadata={"source": filename, "type": "ocr"}))
            except Exception as e:
                logger.warning(f"PDF OCR+Vision skipped: {e}")
        elif ext == ".docx":
            loader = Docx2txtLoader(tmp_path)
            documents = loader.load()
            try:
                from app.services.ocr_helper import ocr_docx_images_v2
                ocr_text = await ocr_docx_images_v2(tmp_path, ai_client=_get_llm_client())
                if ocr_text:
                    documents.append(Document(page_content=ocr_text, metadata={"source": filename, "type": "ocr"}))
            except Exception as e:
                logger.warning(f"DOCX OCR+Vision skipped: {e}")
        elif ext in (".txt", ".md"):
            try:
                text = content.decode("utf-8")
            except UnicodeDecodeError:
                text = content.decode("gbk", errors="ignore")
            documents = [Document(page_content=text, metadata={"source": filename})]

        if not documents:
            raise HTTPException(status_code=400, detail="未能从文件中提取到有效文本")

        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=800,
            chunk_overlap=100,
            separators=["\n\n", "\n", "。", "！", "？", " ", ""],
        )
        chunks = text_splitter.split_documents(documents)

        for i, chunk in enumerate(chunks):
            chunk.metadata["source"] = filename
            chunk.metadata["chunk_index"] = i

        result = await rag_engine.ingest_documents_batch(chunks, "default", user_id=current_user_id)

        logger.info(f"Knowledge uploaded: {filename} → {len(chunks)} chunks in [{result['collection']}]")
        return {
            "status": "success",
            "collection_name": result["collection"],
            "record_count": result["ingested_count"],
            "filename": filename,
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Failed to process knowledge file {filename}: {str(e)}")
        raise HTTPException(status_code=500, detail=f"文档解析失败: {str(e)}")
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)


@router.post("/knowledge/upload-files")
async def upload_knowledge_files(
    files: List[UploadFile] = File(...),
    chunk_size: int = Form(1000),
    chunk_overlap: int = Form(100),
    current_user_id: str = Depends(get_current_user),
):
    """
    多文件、多格式知识库上传接口
    支持 .docx, .pptx, .xlsx, .pdf, .txt, .md, .jsonl
    自动文本分块 (RecursiveCharacterTextSplitter) + 向量化入库
    数据存入用户专属集合 (kb_{user_id})，实现 Collection 物理隔离
    """
    rag_engine = _get_rag_engine()
    if rag_engine is None:
        raise HTTPException(status_code=503, detail="RAG引擎未初始化")

    if not files:
        raise HTTPException(status_code=400, detail="未提供任何文件")

    supported_exts = (".docx", ".pptx", ".xlsx", ".pdf", ".txt", ".md", ".jsonl")
    from langchain_core.documents import Document
    from langchain_text_splitters import RecursiveCharacterTextSplitter

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", "。", "！", "？", "；", " ", ""],
    )

    all_chunks: List[Document] = []
    file_results = []
    import tempfile

    for f in files:
        filename = f.filename or ""
        ext = os.path.splitext(filename)[1].lower()

        if ext not in supported_exts:
            file_results.append({"filename": filename, "status": "skipped", "reason": f"不支持的格式 {ext}"})
            continue

        content = await f.read()

        try:
            documents = []

            if ext == ".jsonl":
                try:
                    text = content.decode("utf-8")
                except UnicodeDecodeError:
                    text = content.decode("gbk", errors="ignore")
                records = []
                for line in text.strip().split("\n"):
                    line = line.strip()
                    if line:
                        try:
                            records.append(json.loads(line))
                        except json.JSONDecodeError:
                            continue
                if records:
                    for idx, record in enumerate(records):
                        field_parts = [f"{k}: {v}" for k, v in record.items() if v is not None and str(v).strip()]
                        page_content = " | ".join(field_parts)
                        documents.append(Document(page_content=page_content, metadata={"source": filename, "record_index": idx}))

            elif ext == ".pdf":
                with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                    tmp.write(content)
                    tmp_path = tmp.name
                try:
                    import pdfplumber
                    with pdfplumber.open(tmp_path) as pdf:
                        for page_num, page in enumerate(pdf.pages):
                            page_text = page.extract_text() or ""
                            if page_text.strip():
                                documents.append(Document(
                                    page_content=page_text,
                                    metadata={"source": filename, "page": page_num + 1},
                                ))
                    try:
                        from app.services.ocr_helper import ocr_pdf_images_v2
                        ocr_text = await ocr_pdf_images_v2(tmp_path, ai_client=_get_llm_client())
                        if ocr_text:
                            documents.append(Document(page_content=ocr_text, metadata={"source": filename, "type": "ocr"}))
                    except Exception as e:
                        logger.warning(f"PDF OCR+Vision skipped for {filename}: {e}")
                finally:
                    if os.path.exists(tmp_path):
                        os.remove(tmp_path)

            elif ext == ".docx":
                with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                    tmp.write(content)
                    tmp_path = tmp.name
                try:
                    from app.services.docx_parser import DocxParser
                    parser = DocxParser()
                    import asyncio as _aio
                    html = await parser.convert_docx_to_html(tmp_path)
                    if html:
                        from bs4 import BeautifulSoup
                        soup = BeautifulSoup(html, "html.parser")
                        text = soup.get_text(separator="\n", strip=True)
                        if text.strip():
                            documents.append(Document(page_content=text, metadata={"source": filename}))
                except Exception:
                    try:
                        from langchain_community.document_loaders import Docx2txtLoader
                        loader = Docx2txtLoader(tmp_path)
                        documents = loader.load()
                    except Exception:
                        pass
                try:
                    from app.services.ocr_helper import ocr_docx_images_v2
                    ocr_text = await ocr_docx_images_v2(tmp_path, ai_client=_get_llm_client())
                    if ocr_text:
                        documents.append(Document(page_content=ocr_text, metadata={"source": filename, "type": "ocr"}))
                except Exception as e:
                    logger.warning(f"DOCX OCR+Vision skipped for {filename}: {e}")
                finally:
                    if os.path.exists(tmp_path):
                        os.remove(tmp_path)

            elif ext == ".pptx":
                with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                    tmp.write(content)
                    tmp_path = tmp.name
                try:
                    from pptx import Presentation as PptxPresentation
                    prs = PptxPresentation(tmp_path)
                    for slide_num, slide in enumerate(prs.slides, 1):
                        slide_parts = []
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    text = para.text.strip()
                                    if text:
                                        slide_parts.append(text)
                        if slide_parts:
                            slide_text = "\n".join(slide_parts)
                            documents.append(Document(
                                page_content=slide_text,
                                metadata={"source": filename, "slide": slide_num},
                            ))
                finally:
                    if os.path.exists(tmp_path):
                        os.remove(tmp_path)

            elif ext == ".xlsx":
                import pandas as pd
                import io
                try:
                    xls = pd.ExcelFile(io.BytesIO(content))
                    for sheet_name in xls.sheet_names:
                        df = pd.read_excel(xls, sheet_name=sheet_name)
                        if not df.empty:
                            text = df.to_string(index=False)
                            documents.append(Document(
                                page_content=text,
                                metadata={"source": filename, "sheet": sheet_name},
                            ))
                except Exception as e:
                    logger.warning(f"XLSX parse failed for {filename}: {e}")

            elif ext in (".txt", ".md"):
                try:
                    text = content.decode("utf-8")
                except UnicodeDecodeError:
                    text = content.decode("gbk", errors="ignore")
                if text.strip():
                    documents.append(Document(page_content=text, metadata={"source": filename}))

            if not documents:
                file_results.append({"filename": filename, "status": "empty", "chunks": 0})
                continue

            chunks = text_splitter.split_documents(documents)
            for i, chunk in enumerate(chunks):
                chunk.metadata["source"] = filename
                chunk.metadata["chunk_index"] = i

            all_chunks.extend(chunks)
            file_results.append({"filename": filename, "status": "success", "chunks": len(chunks)})
            logger.info(f"Knowledge: {filename} → {len(chunks)} chunks")

        except Exception as e:
            logger.error(f"Failed to process {filename}: {e}")
            file_results.append({"filename": filename, "status": "error", "reason": str(e)})

    total_chunks = 0
    effective_collection = ""
    if all_chunks:
        try:
            result = await rag_engine.ingest_documents_batch(all_chunks, "default", user_id=current_user_id)
            total_chunks = result["ingested_count"]
            effective_collection = result["collection"]
            logger.info(f"Knowledge batch: {len(all_chunks)} chunks → [{effective_collection}], ingested={total_chunks}")
        except Exception as e:
            logger.error(f"Knowledge batch ingest failed: {e}")
            raise HTTPException(status_code=500, detail=f"向量化入库失败: {str(e)}")

    try:
        session = get_session()
        for fr in file_results:
            if fr.get("status") == "success":
                doc_meta = DocumentMeta(
                    filename=fr["filename"],
                    file_type=os.path.splitext(fr["filename"])[1].lower().lstrip("."),
                    chunk_count=fr.get("chunks", 0),
                    status="completed",
                    collection_name=effective_collection,
                    owner_id=current_user_id or "",
                    user_id=current_user_id,
                )
                session.add(doc_meta)
        session.commit()
    except Exception as e:
        logger.warning(f"Failed to save document metadata: {e}")
    finally:
        session.close()

    return {
        "status": "success",
        "collection_name": effective_collection,
        "total_chunks": total_chunks,
        "files": file_results,
    }


@router.get("/knowledge/stats")
async def knowledge_stats(
    x_user_id: Optional[str] = Header(None, alias="X-User-Id"),
    current_user_id: Optional[str] = Depends(get_current_user_optional),
):
    effective_user_id = x_user_id or current_user_id or None
    if not effective_user_id:
        raise HTTPException(status_code=401, detail="未提供用户身份")

    rag_engine = _get_rag_engine()
    if rag_engine is None:
        raise HTTPException(status_code=503, detail="RAG引擎未初始化")

    return await rag_engine.get_knowledge_stats(user_id=effective_user_id)


@router.post("/kb/upload")
async def kb_upload(
    files: List[UploadFile] = File(...),
    current_user_id: Optional[str] = Depends(get_current_user_optional),
    x_user_id: Optional[str] = Header(None, alias="X-User-Id"),
):
    """
    专用知识库构建上传接口
    将文件安全存储到用户隔离目录，并触发 Celery 异步任务进行文档向量化和图谱抽取
    """
    effective_user_id = x_user_id or current_user_id or None
    if not effective_user_id:
        raise HTTPException(status_code=401, detail="未提供用户身份，请先登录或携带 X-User-Id")

    supported_exts = (".pdf", ".docx", ".txt", ".md", ".jsonl")
    saved_files = []

    for f in files:
        if not f.filename:
            continue
        ext = os.path.splitext(f.filename)[1].lower()
        if ext not in supported_exts:
            logger.warning(f"KB Upload: skipping unsupported file {f.filename}")
            continue

        file_content = await f.read()
        if len(file_content) > 100 * 1024 * 1024:
            logger.warning(f"KB Upload: file {f.filename} exceeds 100MB limit")
            continue

        kb_dir = os.path.join("uploads", f"user_{effective_user_id}", "kb_docs")
        os.makedirs(kb_dir, exist_ok=True)

        safe_name = f"{uuid.uuid4().hex[:8]}_{f.filename}"
        saved_path = os.path.join(kb_dir, safe_name)
        with open(saved_path, "wb") as out_f:
            out_f.write(file_content)

        saved_files.append({
            "filename": f.filename,
            "saved_path": saved_path,
            "ext": ext,
            "size": len(file_content),
        })
        logger.info(f"KB Upload: saved {f.filename} → {saved_path}")

    if not saved_files:
        raise HTTPException(status_code=400, detail="没有有效的知识库文件被上传")

    from app.models.database import DocumentMeta, get_session
    session = get_session()
    doc_ids = []
    try:
        for sf in saved_files:
            doc = DocumentMeta(
                user_id=effective_user_id,
                filename=sf["filename"],
                file_type=sf["ext"].lstrip("."),
                file_size=sf.get("size", 0),
                status="processing",
                collection_name="default",
                owner_id=effective_user_id,
            )
            session.add(doc)
            session.flush()
            doc_ids.append(doc.id)
            sf["doc_id"] = doc.id
        session.commit()
        logger.info(f"KB Upload: inserted {len(doc_ids)} DocumentMeta records for user {effective_user_id}")
    except Exception as e:
        session.rollback()
        logger.warning(f"KB Upload: failed to insert DocumentMeta: {e}")
    finally:
        session.close()

    from app.core.config import settings
    if settings.USE_CELERY:
        from app.worker.tasks import process_kb_document
        process_kb_document.delay(
            user_id=effective_user_id,
            files=saved_files,
        )
        logger.info(f"KB Upload: {len(saved_files)} files dispatched to Celery for user {effective_user_id}")
    else:
        import asyncio
        from app.services.kb_processor import process_kb_files_sync
        asyncio.create_task(process_kb_files_sync(effective_user_id, saved_files))
        logger.info(f"KB Upload: {len(saved_files)} files processing locally for user {effective_user_id}")

    return {
        "message": "文献已接收，AI 正在后台构建知识网络",
        "status": "processing",
        "file_count": len(saved_files),
        "files": [f["filename"] for f in saved_files],
        "doc_ids": doc_ids,
    }


@router.get("/kb/documents")
async def kb_documents(
    x_user_id: Optional[str] = Header(None, alias="X-User-Id"),
    current_user_id: Optional[str] = Depends(get_current_user_optional),
):
    effective_user_id = x_user_id or current_user_id or None
    if not effective_user_id:
        raise HTTPException(status_code=401, detail="未提供用户身份")

    from app.models.database import DocumentMeta, get_session
    session = get_session()
    try:
        docs = (
            session.query(DocumentMeta)
            .filter(DocumentMeta.user_id == effective_user_id)
            .order_by(DocumentMeta.upload_time.desc())
            .all()
        )
        return {
            "total": len(docs),
            "documents": [doc.to_dict() for doc in docs],
        }
    except Exception as e:
        logger.error(f"KB Documents query failed: {e}")
        return {"total": 0, "documents": []}
    finally:
        session.close()


@router.delete("/kb/documents/{document_id}")
async def delete_kb_document(
    document_id: int,
    x_user_id: Optional[str] = Header(None, alias="X-User-Id"),
    current_user_id: Optional[str] = Depends(get_current_user_optional),
):
    effective_user_id = x_user_id or current_user_id or None
    if not effective_user_id:
        raise HTTPException(status_code=401, detail="未提供用户身份")

    from app.models.database import DocumentMeta, get_session
    session = get_session()
    try:
        doc = session.query(DocumentMeta).filter(
            DocumentMeta.id == document_id,
            DocumentMeta.user_id == effective_user_id,
        ).first()

        if not doc:
            raise HTTPException(status_code=404, detail="文档不存在或无权删除")

        filename = doc.filename
        file_type = doc.file_type

        rag_engine = _get_rag_engine()
        if rag_engine:
            try:
                deleted_chunks = await rag_engine.delete_documents_by_source(filename, effective_user_id)
                logger.info(f"KB Delete: removed {deleted_chunks} vector chunks for '{filename}'")
            except Exception as e:
                logger.warning(f"KB Delete: vector cleanup failed for '{filename}': {e}")

        kb_dir = os.path.join("uploads", f"user_{effective_user_id}", "kb_docs")
        if os.path.isdir(kb_dir):
            for f in os.listdir(kb_dir):
                if f.endswith(f"_{filename}") or f == filename:
                    file_path = os.path.join(kb_dir, f)
                    try:
                        if os.path.exists(file_path):
                            os.remove(file_path)
                            logger.info(f"KB Delete: removed physical file {file_path}")
                    except Exception as e:
                        logger.warning(f"KB Delete: failed to remove file {file_path}: {e}")

        session.delete(doc)
        session.commit()
        logger.info(f"KB Delete: deleted DocumentMeta id={document_id} filename='{filename}' for user {effective_user_id}")

        return {"status": "success", "message": "文献及其 AI 记忆已彻底清除", "deleted_document": filename}
    except HTTPException:
        raise
    except Exception as e:
        session.rollback()
        logger.error(f"KB Delete failed: {e}")
        raise HTTPException(status_code=500, detail=f"删除失败: {str(e)}")
    finally:
        session.close()


@router.get("/flashcards/due", deprecated=True)
async def flashcards_due(
    x_user_id: Optional[str] = Header(None, alias="X-User-Id"),
    current_user_id: Optional[str] = Depends(get_current_user_optional),
):
    raise HTTPException(status_code=410, detail="此接口已废弃，请使用 GET /flashcards/draw")


@router.post("/flashcards/{card_id}/review", deprecated=True)
async def review_flashcard(
    card_id: int,
    review: ReviewSubmit,
    x_user_id: Optional[str] = Header(None, alias="X-User-Id"),
    current_user_id: Optional[str] = Depends(get_current_user_optional),
):
    raise HTTPException(status_code=410, detail="此接口已废弃，闪卡系统已改为实时生成模式")


@router.get("/flashcards/draw")
async def flashcards_draw(
    limit: int = Query(5, ge=1, le=20, description="生成卡片数量"),
    x_user_id: str = Header(..., alias="X-User-Id"),
):
    import os
    import random
    import re
    from datetime import datetime, timedelta
    from app.models.database import ChunkExposureHistory, chunk_text_hash, get_session

    rag_engine = _get_rag_engine()
    llm_client = _get_llm_client()

    if rag_engine is None or llm_client is None:
        logger.warning("flashcards_draw: RAG engine or LLM client not available")
        return {"status": "success", "total": 0, "cards": []}

    try:
        from app.models.database import DocumentMeta

        db_sess = get_session()
        try:
            user_docs = db_sess.query(DocumentMeta).filter(
                DocumentMeta.user_id == x_user_id,
                DocumentMeta.status == "completed",
            ).all()
        finally:
            db_sess.close()

        if not user_docs:
            return {
                "status": "success",
                "total": 0,
                "cards": [],
                "reason": "no_documents",
                "message": "请先上传文档至知识库",
            }

        selected_doc = random.choice(user_docs)
        doc_filename = selected_doc.filename
        doc_name = os.path.splitext(doc_filename)[0]

        noise_words = [
            "死锁", "并发", "内存管理", "TCP/IP", "路由",
            "时间复杂度", "二叉树", "图论", "缓存", "进程调度",
            "哈希", "排序", "总线", "指令集", "页表",
        ]
        noise_sample = random.sample(noise_words, 2)
        dynamic_query = f"{doc_name} {' '.join(noise_sample)} 原理 机制 核心"

        chunks = await rag_engine.mmr_search(
            query=dynamic_query,
            user_id=x_user_id,
            fetch_k=40,
            k=20,
            lambda_mult=0.5,
        )

        if not chunks:
            logger.info(f"flashcards_draw: no chunks for query='{dynamic_query}', user={x_user_id}")
            return {"status": "success", "total": 0, "cards": []}

        logger.info(f"flashcards_draw: global MMR recalled {len(chunks)} candidates for query='{dynamic_query}'")

        for chunk in chunks:
            text = chunk.get("text", chunk.get("page_content", chunk.get("content", "")))
            chunk["_hash"] = chunk_text_hash(text)

        hash_list = [c["_hash"] for c in chunks]

        session = get_session()
        try:
            history_records = session.query(ChunkExposureHistory).filter(
                ChunkExposureHistory.user_id == x_user_id,
                ChunkExposureHistory.chunk_text_hash.in_(hash_list),
            ).all()

            history_map = {}
            for rec in history_records:
                history_map[rec.chunk_text_hash] = rec
        finally:
            session.close()

        now = datetime.utcnow()

        approved_chunks = []
        rejected_chunks = []

        for chunk in chunks:
            h = chunk["_hash"]
            if h not in history_map:
                approved_chunks.append((chunk, 999999.0))
            else:
                rec = history_map[h]
                delta_hours = (now - rec.last_exposed_at).total_seconds() / 3600.0
                if rec.consecutive_correct == 0:
                    score = 888888.0
                else:
                    score = delta_hours - (rec.consecutive_correct * 24.0)
                if score >= 0.5:
                    approved_chunks.append((chunk, score))
                else:
                    rejected_chunks.append((chunk, score))

        approved_chunks.sort(key=lambda x: -x[1])
        rejected_chunks.sort(key=lambda x: x[1])

        card_mode = "new"

        if len(approved_chunks) >= 5:
            selected = [item[0] for item in approved_chunks[:5]]
        elif len(approved_chunks) > 0:
            need = 5 - len(approved_chunks)
            selected = [item[0] for item in approved_chunks]
            selected += [item[0] for item in rejected_chunks[:need]]
            card_mode = "review"
        else:
            if rejected_chunks:
                selected = [item[0] for item in rejected_chunks[:5]]
                card_mode = "review"
            else:
                logger.info(f"flashcards_draw: no chunks at all after scoring, user={x_user_id}")
                return {"status": "success", "total": 0, "cards": [], "reason": "all_recently_seen"}

        top_approved_str = f"top_approved={approved_chunks[0][1]:.2f}" if approved_chunks else "top_approved=N/A"
        logger.info(
            f"flashcards_draw: elastic rerank → "
            f"approved={len(approved_chunks)}, rejected={len(rejected_chunks)}, "
            f"selected={len(selected)}, mode={card_mode}, {top_approved_str}"
        )

        chunk_hash_map = {}
        for chunk in selected:
            h = chunk["_hash"]
            source = chunk.get("metadata", {}).get("source", "未知文档")
            chunk_hash_map[source] = h

        context_parts = []
        for i, chunk in enumerate(selected):
            text = chunk.get("text", chunk.get("page_content", chunk.get("content", "")))
            source = chunk.get("metadata", {}).get("source", "未知文档")
            context_parts.append(f"[片段{i+1}] (来源: {source})\n{text}")

        context = "\n\n".join(context_parts)

        if card_mode == "review":
            prompt = f"""你是一个严苛的考试出题专家。以下是用户近期学过但可能遗忘的知识片段，请生成 {limit} 张【深度复习题】闪卡。

要求：
- 问题角度要比之前更刁钻，侧重易混淆点、边界条件、底层原理
- 答案要精准，不超过 60 字
- 优先考察：概念辨析、原理推导、常见误区、对比差异
- 每张卡必须包含 document_name 字段，值为来源文件名
- 必须严格输出 JSON 格式：{{"cards": [{{"question": "...", "answer": "...", "document_name": "来源文件名"}}]}}

参考文本（用户已学过，需要深度复习）：
{context}"""
        else:
            prompt = f"""你是一个严苛的考试出题专家。请阅读以下文本，提取出最核心的知识点，生成 {limit} 张 Q&A 形式的复习闪卡。

要求：
- 问题要简明扼要，直击核心
- 答案要精准，不超过 50 字
- 优先提取定义、原理、公式、关键区别等高价值考点
- 每张卡必须包含 document_name 字段，值为来源文件名
- 必须严格输出 JSON 格式：{{"cards": [{{"question": "...", "answer": "...", "document_name": "来源文件名"}}]}}

参考文本：
{context}"""

        messages = [
            {"role": "system", "content": "你是一个无情的 JSON 生成机器。你必须且只能输出合法的 JSON 格式。绝对禁止使用 Markdown 格式（不要输出 ```json），绝对禁止输出任何寒暄、解释或前缀后缀！只输出纯 JSON！"},
            {"role": "user", "content": prompt},
        ]

        response = await llm_client.acall_api(messages, max_tokens=2048)

        if not response or not response.strip():
            logger.warning("flashcards_draw: LLM returned empty")
            return {"status": "success", "total": 0, "cards": []}

        logger.info(f"flashcards_draw: LLM response (first 300 chars): {response[:300]}")

        try:
            parsed = clean_and_parse_json(response)
        except (ValueError, json.JSONDecodeError) as parse_err:
            logger.warning(f"flashcards_draw: clean_and_parse_json failed: {parse_err}")
            return {"status": "success", "total": 0, "cards": []}

        cards = parsed.get("cards", [])
        if not cards:
            cleaned = re.sub(r'<think[^>]*>.*?</think\s*>', '', response, flags=re.DOTALL)
            pairs = re.findall(r'"question"\s*:\s*"(.+?)"\s*,\s*"answer"\s*:\s*"(.+?)"', cleaned)
            if pairs:
                cards = [{"question": q, "answer": a, "document_name": "未知文档"} for q, a in pairs]

        valid_cards = []
        for card in cards:
            q = card.get("question", "").strip()
            a = card.get("answer", "").strip()
            doc = card.get("document_name", "").strip() or "未知文档"
            if q and a:
                valid_cards.append({
                    "question": q,
                    "answer": a,
                    "document_name": doc,
                    "chunk_hash": chunk_hash_map.get(doc, ""),
                })

        logger.info(f"flashcards_draw: generated {len(valid_cards)} cards for user={x_user_id}, mode={card_mode}")
        return {"status": "success", "total": len(valid_cards), "cards": valid_cards, "mode": card_mode}

    except Exception as e:
        logger.error(f"flashcards_draw failed: {e}", exc_info=True)
        return {"status": "success", "total": 0, "cards": []}


@router.post("/flashcards/review")
async def flashcards_review(
    chunk_hash: str = Body(..., embed=True),
    is_correct: bool = Body(..., embed=True),
    x_user_id: str = Header(..., alias="X-User-Id"),
):
    from datetime import datetime
    from app.models.database import ChunkExposureHistory, get_session

    if not chunk_hash or not chunk_hash.strip():
        raise HTTPException(status_code=400, detail="chunk_hash 不能为空")

    session = get_session()
    try:
        rec = session.query(ChunkExposureHistory).filter(
            ChunkExposureHistory.user_id == x_user_id,
            ChunkExposureHistory.chunk_text_hash == chunk_hash,
        ).first()

        now = datetime.utcnow()

        if rec:
            rec.last_exposed_at = now
            rec.exposure_count += 1
            if is_correct:
                rec.consecutive_correct += 1
            else:
                rec.consecutive_correct = 0
        else:
            rec = ChunkExposureHistory(
                user_id=x_user_id,
                chunk_text_hash=chunk_hash,
                last_exposed_at=now,
                exposure_count=1,
                consecutive_correct=1 if is_correct else 0,
            )
            session.add(rec)

        session.commit()
        logger.info(
            f"flashcards_review: user={x_user_id}, hash={chunk_hash[:8]}..., "
            f"is_correct={is_correct}, consecutive_correct={rec.consecutive_correct}"
        )
        return {
            "status": "success",
            "consecutive_correct": rec.consecutive_correct,
            "exposure_count": rec.exposure_count,
        }
    except Exception as e:
        session.rollback()
        logger.error(f"flashcards_review failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        session.close()


@router.get("/task/{task_id}", response_model=TaskProgress)
async def get_task_status(task_id: str):
    if task_id not in _tasks:
        raise HTTPException(status_code=404, detail=f"任务 {task_id} 不存在")

    task = _tasks[task_id]
    return TaskProgress(
        task_id=task["task_id"],
        status=task["status"],
        progress=task["progress"],
        stage=task["stage"],
        tables_total=task["tables_total"],
        tables_completed=task["tables_completed"],
        error_message=task["error_message"],
        created_at=task["created_at"],
        updated_at=task["updated_at"],
    )


_redis_for_sse = None


def _get_redis_for_sse():
    global _redis_for_sse
    if _redis_for_sse is None:
        from app.core.celery_app import REDIS_PASSWORD, REDIS_HOST, REDIS_PORT
        if REDIS_PASSWORD:
            url = f"redis://:{REDIS_PASSWORD}@{REDIS_HOST}:{REDIS_PORT}/0"
        else:
            url = f"redis://{REDIS_HOST}:{REDIS_PORT}/0"
        _redis_for_sse = redis_lib.from_url(url, decode_responses=True)
    return _redis_for_sse


@router.post("/task/submit")
async def submit_task(
    background_tasks: BackgroundTasks,
    files: Optional[List[UploadFile]] = File(None),
    user_instruction: str = Form(""),
    model: Optional[str] = Form(None),
    max_retries: int = Form(3),
    current_user_id: str = Depends(get_current_user),
    x_user_id: Optional[str] = Header(None, alias="X-User-Id"),
):
    """
    提交文档处理任务，立即返回 task_id。
    根据 USE_CELERY 配置自动选择调度引擎:
      - USE_CELERY=True:  提交到 Celery 队列 (生产环境)
      - USE_CELERY=False: FastAPI BackgroundTasks 本地执行 (开发环境)
    前端拿到 task_id 后，通过 GET /v1/task/{task_id}/stream 监听进度。
    """
    from app.core.config import settings
    from app.services.rag_engine import RAGEngine

    task_id = uuid.uuid4().hex

    uploaded_files_info = []
    if files:
        from app.core.file_manager import file_manager
        effective_user_id = x_user_id or current_user_id or None
        for f in files:
            if not f.filename:
                continue
            file_type_val = SupportedFileType.from_filename(f.filename)
            if file_type_val is None:
                continue
            file_content = await f.read()
            if len(file_content) > 200 * 1024 * 1024:
                raise HTTPException(status_code=413, detail=f"文件 {f.filename} 大小超过200MB限制")
            path = file_manager.save_upload(f.filename, file_content, file_type_val.value, user_id=effective_user_id)
            uploaded_files_info.append({
                "path": path,
                "type": file_type_val.value,
                "filename": f.filename,
            })

    doc_ids = []
    db_session = get_session()
    try:
        for f_info in uploaded_files_info:
            doc_meta = DocumentMeta(
                user_id=current_user_id,
                filename=f_info.get("filename", "unknown"),
                file_type=f_info.get("type", "unknown"),
                file_size=0,
                status="processing",
                task_id=task_id,
                owner_id=current_user_id,
                collection_name=RAGEngine._make_collection_name(current_user_id),
            )
            db_session.add(doc_meta)
        db_session.commit()
        for f_info in uploaded_files_info:
            doc = db_session.query(DocumentMeta).filter(
                DocumentMeta.task_id == task_id,
                DocumentMeta.filename == f_info.get("filename", "unknown"),
            ).first()
            if doc:
                doc_ids.append(doc.id)
    except Exception as e:
        logger.warning(f"Failed to create DocumentMeta records: {e}")
        db_session.rollback()
    finally:
        db_session.close()

    if settings.USE_CELERY:
        from app.worker.tasks import process_document_task
        process_document_task.delay(
            task_id=task_id,
            uploaded_files=uploaded_files_info,
            user_instruction=user_instruction,
            max_retries=max_retries,
            thread_id=f"thread_{uuid.uuid4().hex[:8]}",
            user_id=current_user_id,
            doc_ids=doc_ids,
        )
        logger.info(f"Task submitted [Celery]: {task_id}, files={len(uploaded_files_info)}")
    else:
        from app.core.task_runner import run_document_pipeline_sync
        background_tasks.add_task(
            run_document_pipeline_sync,
            task_id=task_id,
            uploaded_files=uploaded_files_info,
            user_instruction=user_instruction,
            max_retries=max_retries,
            thread_id=f"thread_{uuid.uuid4().hex[:8]}",
            user_id=current_user_id,
            doc_ids=doc_ids,
            use_redis=False,
        )
        logger.info(f"Task submitted [Local]: {task_id}, files={len(uploaded_files_info)}")

    if current_user_id and user_instruction:
        try:
            from app.services.memory_service import memory_manager
            background_tasks.add_task(
                _bg_extract_memory,
                user_instruction=user_instruction,
                user_id=current_user_id,
                llm_client=_get_llm_client(),
                rag_engine=_get_rag_engine(),
            )
        except Exception as e:
            logger.warning(f"Memory background task scheduling failed: {e}")

    return {"task_id": task_id, "status": "processing", "engine": "celery" if settings.USE_CELERY else "local"}


@router.get("/task/{task_id}/stream")
async def task_stream(task_id: str, current_user_id: str = Depends(get_current_user_from_header_or_query)):
    """
    SSE 接口：实时推送任务进度更新。
    根据 USE_CELERY 配置自动选择进度源:
      - Celery 模式: Redis Pub/Sub 监听
      - 本地模式:   轮询 ProgressStore 内存 + SQLite DocumentMeta 状态
    """
    from app.core.config import settings

    if settings.USE_CELERY:
        return await _task_stream_redis(task_id)
    else:
        return await _task_stream_local(task_id)


async def _task_stream_redis(task_id: str):
    r = _get_redis_for_sse()
    pubsub = r.pubsub()
    channel = f"task_channel:{task_id}"
    pubsub.subscribe(channel)

    async def event_generator():
        last_progress = None
        idle_count = 0
        max_idle = 600

        while True:
            cached = r.get(f"task_progress:{task_id}")
            if cached:
                try:
                    payload = json.loads(cached)
                    progress_key = (payload.get("progress"), payload.get("action"))
                    if progress_key != last_progress:
                        last_progress = progress_key
                        idle_count = 0
                        sse_data = json.dumps(payload, ensure_ascii=False)
                        yield f"data: {sse_data}\n\n"

                        if payload.get("status") in ("success", "error"):
                            result = payload.get("result") or payload.get("error")
                            yield f"data: {json.dumps({'status': payload['status'], 'result': result}, ensure_ascii=False, default=str)}\n\n"
                            break
                except json.JSONDecodeError:
                    pass

            message = pubsub.get_message(timeout=1.0)
            if message and message["type"] == "message":
                try:
                    payload = json.loads(message["data"])
                    progress_key = (payload.get("progress"), payload.get("action"))
                    if progress_key != last_progress:
                        last_progress = progress_key
                        idle_count = 0
                        sse_data = json.dumps(payload, ensure_ascii=False)
                        yield f"data: {sse_data}\n\n"

                        if payload.get("status") in ("success", "error"):
                            result = payload.get("result") or payload.get("error")
                            yield f"data: {json.dumps({'status': payload['status'], 'result': result}, ensure_ascii=False, default=str)}\n\n"
                            break
                except json.JSONDecodeError:
                    pass

            idle_count += 1
            if idle_count >= max_idle:
                yield f"data: {json.dumps({'status': 'timeout', 'action': '等待超时'})}\n\n"
                break

            await asyncio.sleep(0.5)

        pubsub.unsubscribe(channel)
        pubsub.close()

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


async def _task_stream_local(task_id: str):
    from app.core.task_runner import progress_store

    async def event_generator():
        last_progress = None
        idle_count = 0
        max_idle = 600

        while True:
            payload = progress_store.get(task_id)

            if payload:
                progress_key = (payload.get("progress"), payload.get("action"))
                if progress_key != last_progress:
                    last_progress = progress_key
                    idle_count = 0
                    sse_data = json.dumps(payload, ensure_ascii=False)
                    yield f"data: {sse_data}\n\n"

                    if payload.get("status") in ("success", "error"):
                        result = payload.get("result") or payload.get("error")
                        yield f"data: {json.dumps({'status': payload['status'], 'result': result}, ensure_ascii=False, default=str)}\n\n"
                        break
            else:
                try:
                    session = get_session()
                    doc = session.query(DocumentMeta).filter(
                        DocumentMeta.task_id == task_id,
                        DocumentMeta.user_id == current_user_id,
                    ).first()
                    if doc and doc.status in ("completed", "failed"):
                        status_payload = {
                            "task_id": task_id,
                            "progress": 100 if doc.status == "completed" else -1,
                            "action": "处理完成" if doc.status == "completed" else "处理失败",
                            "status": "success" if doc.status == "completed" else "error",
                        }
                        sse_data = json.dumps(status_payload, ensure_ascii=False)
                        yield f"data: {sse_data}\n\n"
                        session.close()
                        break
                    session.close()
                except Exception:
                    pass

            idle_count += 1
            if idle_count >= max_idle:
                yield f"data: {json.dumps({'status': 'timeout', 'action': '等待超时'})}\n\n"
                break

            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )
