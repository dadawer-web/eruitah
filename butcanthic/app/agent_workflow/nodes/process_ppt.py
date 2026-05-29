import logging
from typing import Any, Dict

from app.agent_workflow.state import WorkflowState
from app.services.ppt_service import PPTAnalyzer, PPTGenerator

logger = logging.getLogger(__name__)


async def process_ppt_node(
    state: WorkflowState,
    ppt_analyzer: PPTAnalyzer,
    llm_client=None,
    rag_engine=None,
) -> WorkflowState:
    file_path = state.get("file_path", "")
    user_instruction = state.get("user_instruction", "")
    user_id = state.get("user_id", "")

    logger.info(f"ProcessPPT: starting for {file_path}")

    rag_context = ""
    if rag_engine and user_id:
        try:
            query = user_instruction if user_instruction else "演示文稿内容"
            results = await rag_engine.semantic_search(query=query, top_k=3, user_id=user_id)
            if results:
                rag_parts = []
                for r in results[:3]:
                    content = r.get("content", "")
                    if content:
                        rag_parts.append(content[:500])
                if rag_parts:
                    rag_context = "\n\n".join(rag_parts)
                    logger.info(f"ProcessPPT: RAG retrieved {len(rag_parts)} knowledge chunks for context")
        except Exception as e:
            logger.warning(f"ProcessPPT: RAG retrieval failed: {e}")

    try:
        result = await ppt_analyzer.extract_text_and_structure(file_path)
    except Exception as e:
        logger.warning(f"ProcessPPT: PPTAnalyzer failed: {e}, using safe fallback")
        try:
            import pptx
            prs = pptx.Presentation(file_path)
            text_chunks = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_chunks.append(shape.text)
            fallback_text = "\n".join(text_chunks)
            result = {
                "full_text": fallback_text,
                "slide_count": len(prs.slides),
                "slides": [{"index": i, "title": f"Slide {i + 1}", "texts": [], "has_table": False, "has_image": False, "layout_name": ""} for i in range(len(prs.slides))],
            }
        except Exception as inner_e:
            logger.error(f"ProcessPPT: Safe fallback completely failed: {inner_e}")
            return {
                **state,
                "error_message": f"PPT分析失败: {e}",
            }

    structured_data = {
        "slide_count": result["slide_count"],
        "slides": result["slides"],
        "full_text": result["full_text"],
    }

    logger.info(f"ProcessPPT: {result['slide_count']} slides extracted, generating SlideDeck JSON...")

    try:
        generator = PPTGenerator(llm_client)
        slide_deck = await generator.generate_slide_deck(
            user_instruction=user_instruction or "请根据PPT内容生成演示文稿",
            ppt_file_path=file_path,
            rag_context=rag_context,
        )
        structured_data["slide_deck"] = slide_deck
        logger.info(f"ProcessPPT: SlideDeck JSON generated with {len(slide_deck.get('slides', []))} slides")
    except Exception as e:
        logger.error(f"ProcessPPT: SlideDeck generation failed: {repr(e)}")

    return {
        **state,
        "structured_data": structured_data,
        "error_message": "",
    }
