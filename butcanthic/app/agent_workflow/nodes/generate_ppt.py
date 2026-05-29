import json
import logging
import re
import urllib.parse
from datetime import datetime
from typing import Any, Dict, Optional

from app.agent_workflow.state import WorkflowState
from app.core.prompt_manager import PromptManager
from app.models.ppt_schema import Presentation

logger = logging.getLogger(__name__)


async def generate_ppt_node(
    state: WorkflowState,
    llm_client=None,
    rag_engine=None,
) -> WorkflowState:
    user_instruction = state.get("user_instruction", "")
    file_path = state.get("file_path", "")
    file_type = state.get("file_type", "")
    feedback = state.get("feedback", "")
    retry_count = state.get("retry_count", 0)

    logger.info(f"GeneratePPT: starting, instruction='{user_instruction[:100]}', retry={retry_count}, has_feedback={bool(feedback)}")

    uploaded_files = state.get("uploaded_files", [])
    document_content = ""
    collected_images = []

    if uploaded_files:
        for f in uploaded_files:
            try:
                result = await _extract_document_content(f.get("path"), f.get("type"))
                filename = f.get("filename", "未知文件")
                text = result.get("text", "")
                images = result.get("images", [])
                if text:
                    document_content += f"\n\n=== 【参考资料来源：{filename}】 ===\n{text}"
                if images:
                    collected_images.extend(images[:5])
            except Exception as e:
                logger.error(f"提取文件 {f.get('filename')} 失败: {e}")
        document_content = document_content.strip()
    else:
        result = await _extract_document_content(file_path, file_type)
        document_content = result.get("text", "")
        collected_images = result.get("images", [])

    existing_ppt = state.get("ppt_data", {})
    history_context = ""
    if existing_ppt and existing_ppt.get("slides"):
        existing_json = json.dumps(existing_ppt, ensure_ascii=False)
        if len(existing_json) > 15000:
            summary_slides = []
            for s in existing_ppt["slides"][:8]:
                summary_slides.append({
                    "layout": s.get("layout", ""),
                    "title": s.get("title", ""),
                    "subtitle": s.get("subtitle", ""),
                    "components_count": len(s.get("components", [])),
                })
            existing_json = json.dumps({"meta": existing_ppt.get("meta", {}), "design": existing_ppt.get("design", {}), "slides_summary": summary_slides, "total_slides": len(existing_ppt["slides"])}, ensure_ascii=False)
        history_context = (
            f"\n\n【⚠️重要：这是上一轮已经生成的 PPT 底稿】\n{existing_json}\n\n"
            "【增量修改指令】：用户的最新指令是针对这份已有 PPT 进行的修改或扩充。"
            "请你在此基础上进行增删改，必须返回**包含所有幻灯片页面的完整新 PPT JSON**，"
            "绝对不要脱离原有内容重新生成一份毫不相关的！"
        )
        logger.info(f"GeneratePPT: detected existing PPT with {len(existing_ppt['slides'])} slides, enabling incremental mode")

    feedback_context = ""
    if feedback:
        feedback_context = (
            f"\n\n【⚠️重要：这是你上次生成的错误结果，请根据以下反馈进行修正】\n"
            f"审查反馈：{feedback}\n\n"
            "请务必针对上述问题逐一修正，生成符合要求的新 PPT JSON。"
        )
        logger.info(f"GeneratePPT: injected critic feedback (len={len(feedback)})")

    knowledge_context = ""
    messages = state.get("messages", [])
    if messages:
        research_parts = []
        for msg in messages:
            if hasattr(msg, 'name') and msg.name in ("Web_Researcher", "Knowledge_Librarian") and msg.content:
                research_parts.append(msg.content)
        if research_parts:
            knowledge_context = "\n\n".join(research_parts)
            logger.info(f"GeneratePPT: injected {len(research_parts)} research results from messages")

    cross_agent_context = ""
    structured_data = state.get("structured_data", {})
    if structured_data:
        excel_output = ""
        if "output" in structured_data:
            excel_output = str(structured_data.get("output", ""))
        elif "log" in structured_data and isinstance(structured_data.get("log"), list):
            excel_output = "\n".join(str(l) for l in structured_data["log"])

        if structured_data.get("schema") or excel_output:
            if excel_output:
                cross_agent_context = (
                    f"\n\n【⚠️ 核心数据支撑】：这是 DataAgent 刚刚通过 Python 代码对源数据进行硬核分析得出的结论，"
                    f"请务必将其作为核心论点写入 PPT 中：\n{excel_output[:5000]}"
                )
            else:
                schema_summary = json.dumps(structured_data, ensure_ascii=False, default=str)[:3000]
                cross_agent_context = (
                    f"\n\n【⚠️ 数据分析结果】：这是 DataAgent 对 Excel 数据的分析结论，请将其融入 PPT：\n{schema_summary}"
                )
            logger.info("GeneratePPT: Injected Excel DataAgent conclusions into context")

    if rag_engine and not knowledge_context:
        try:
            user_id = state.get("user_id", "")
            if user_id:
                queries = _extract_search_queries(user_instruction, document_content)
                results = []
                for q in queries:
                    docs = await rag_engine.semantic_search(q, top_k=3, user_id=user_id)
                    for doc in docs:
                        text = doc.get("text", doc.get("page_content", ""))
                        if text:
                            results.append(text)
                if results:
                    knowledge_context = "\n".join(results[:10])
                    logger.info(f"GeneratePPT: retrieved {len(results)} knowledge chunks from kb_user_{user_id[:8]}...")
            else:
                logger.warning("GeneratePPT: user_id 缺失，跳过知识库检索（Collection 物理隔离要求必须提供 user_id）")
        except Exception as e:
            logger.warning(f"GeneratePPT: knowledge retrieval failed: {e}")

    presentation_data = None

    if llm_client and hasattr(llm_client, 'langchain_llm') and llm_client.langchain_llm:
        try:
            presentation_data = await _generate_with_structured_output(
                llm_client, user_instruction, document_content, knowledge_context, history_context, cross_agent_context, feedback_context, collected_images
            )
        except Exception as e:
            logger.error(f"🚨 [GeneratePPT] with_structured_output failed: {repr(e)}")
            presentation_data = None

    if presentation_data is None and llm_client:
        try:
            presentation_data = await _generate_with_call_api(
                llm_client, user_instruction, document_content, knowledge_context, history_context, cross_agent_context, feedback_context, collected_images
            )
        except Exception as e:
            logger.error(f"🚨 [GeneratePPT] call_api fallback failed: {repr(e)}")
            presentation_data = None

    if presentation_data is None:
        logger.warning("GeneratePPT: all LLM methods failed, using fallback")
        presentation_data = _build_fallback_presentation(user_instruction, document_content)

    await _inject_image_urls(presentation_data)

    current_structured = state.get("structured_data", {})
    updated_structured = {**current_structured, "presentation": presentation_data}

    slide_count = len(presentation_data.get('slides', []))
    logger.info(
        f"GeneratePPT: completed, {slide_count} slides generated"
    )

    return {
        **state,
        "structured_data": updated_structured,
        "ppt_data": presentation_data,
        "error_message": "",
        "current_progress": 85,
        "current_action": f"PPT生成完毕，共 {slide_count} 页幻灯片",
    }


async def _generate_with_structured_output(
    llm_client,
    user_instruction: str,
    document_content: str,
    knowledge_context: str,
    history_context: str = "",
    cross_agent_context: str = "",
    feedback_context: str = "",
    images: list = None,
) -> Optional[Dict[str, Any]]:
    from langchain_core.messages import HumanMessage, SystemMessage

    logger.info("🚀 [GeneratePPT] using with_structured_output(Presentation)")

    structured_llm = llm_client.langchain_llm.with_structured_output(Presentation)

    user_text = PromptManager.get_prompt(
        "generate_ppt_user",
        user_instruction=user_instruction or "请提炼文档核心内容生成PPT",
        document_content=document_content[:250000] if document_content else "（无文档内容）",
        knowledge_context=knowledge_context[:5000] if knowledge_context else "（无知识库补充）",
        cross_agent_context=cross_agent_context[:5000] if cross_agent_context else "",
        history_context=history_context,
        feedback_context=feedback_context,
    )

    if images:
        logger.info(f"GeneratePPT: injecting {len(images)} images into structured_output request")

    user_content = user_text
    if images:
        user_content = [{"type": "text", "text": user_text}]
        for img in images[:5]:
            user_content.append({"type": "image_url", "image_url": {"url": img}})

    messages = [
        SystemMessage(content=PromptManager.get_prompt("generate_ppt_system")),
        HumanMessage(content=user_content),
    ]

    result = await structured_llm.ainvoke(messages, config={"max_tokens": 8192})

    if result and isinstance(result, Presentation):
        logger.info(f"📋 [GeneratePPT] with_structured_output success: {len(result.slides)} slides")
        return result.model_dump()
    elif result and isinstance(result, dict):
        if "presentation" in result and "slides" not in result:
            result = result["presentation"]
        result = _sanitize_presentation_dict(result)
        logger.info(f"📋 [GeneratePPT] with_structured_output returned dict")
        return result

    logger.warning("🚨 [GeneratePPT] with_structured_output returned None")
    return None


async def _generate_with_call_api(
    llm_client,
    user_instruction: str,
    document_content: str,
    knowledge_context: str,
    history_context: str = "",
    cross_agent_context: str = "",
    feedback_context: str = "",
    images: list = None,
) -> Optional[Dict[str, Any]]:
    logger.info("🚀 [GeneratePPT] using call_api fallback")

    user_text = PromptManager.get_prompt(
        "generate_ppt_user",
        user_instruction=user_instruction or "请提炼文档核心内容生成PPT",
        document_content=document_content[:250000] if document_content else "（无文档内容）",
        knowledge_context=knowledge_context[:5000] if knowledge_context else "（无知识库补充）",
        cross_agent_context=cross_agent_context[:5000] if cross_agent_context else "",
        history_context=history_context,
        feedback_context=feedback_context,
    )

    user_content = user_text
    if images:
        logger.info(f"GeneratePPT: injecting {len(images)} images into call_api request via 'images' field")

    messages = [
        {"role": "system", "content": PromptManager.get_prompt("generate_ppt_system")},
        {"role": "user", "content": user_content, **({"images": images[:5]} if images else {})},
    ]

    response = await llm_client.acall_api(messages, max_tokens=8192)
    if not response:
        logger.error("🚨 [GeneratePPT] call_api returned empty response")
        return None

    logger.info(f"📋 [GeneratePPT] call_api response length: {len(response)}")
    return _parse_presentation_json(response)


def _sanitize_presentation_dict(parsed: dict) -> dict:
    if not isinstance(parsed, dict):
        return parsed
    if "slides" in parsed and isinstance(parsed["slides"], list):
        for slide in parsed["slides"]:
            if "components" in slide and isinstance(slide["components"], list):
                for comp in slide["components"]:
                    if isinstance(comp, dict):
                        t = comp.get("type", "")
                        if t in ["title", "subtitle", "eyebrow", "paragraph", "description"]:
                            comp["type"] = "heading" if t == "title" else "text"

                        if "items" in comp and isinstance(comp["items"], list):
                            new_items = []
                            for item in comp["items"]:
                                if isinstance(item, dict):
                                    vals = [str(v) for v in item.values() if v]
                                    new_items.append(" - ".join(vals))
                                else:
                                    new_items.append(str(item))
                            comp["items"] = new_items
    return parsed


def _parse_presentation_json(response: str) -> Optional[Dict[str, Any]]:
    if not response or not response.strip():
        return None

    cleaned = response.strip()

    code_block = re.search(r"```(?:json)?\s*\n?(.*?)\n?\s*```", cleaned, re.DOTALL)
    if code_block:
        cleaned = code_block.group(1).strip()

    brace_match = re.search(r"\{.*\}", cleaned, re.DOTALL)
    if brace_match:
        cleaned = brace_match.group(0)

    parsed = None

    try:
        parsed = json.loads(cleaned)
    except Exception as e:
        logger.warning(f"🚨 [GeneratePPT] JSON parse primary failed: {e}")
        try:
            fixed = re.sub(r"'", '"', cleaned)
            fixed = re.sub(r',\s*}', '}', fixed)
            fixed = re.sub(r',\s*]', ']', fixed)
            parsed = json.loads(fixed)
        except Exception as e2:
            logger.warning(f"🚨 [GeneratePPT] JSON parse secondary failed: {e2}")

    if not parsed or not isinstance(parsed, dict):
        logger.warning("🚨 [GeneratePPT] JSON parse failed completely")
        return None

    if "presentation" in parsed and "slides" not in parsed:
        parsed = parsed["presentation"]

    if not isinstance(parsed, dict):
        logger.warning("🚨 [GeneratePPT] Unwrapped data is not a dict")
        return None

    parsed = _sanitize_presentation_dict(parsed)

    try:
        deck = Presentation(**parsed)
        return deck.model_dump()
    except Exception as e:
        logger.warning(f"🚨 [GeneratePPT] Pydantic validation failed: {e}, forcing raw dict with defaults")

    if "slides" not in parsed or not isinstance(parsed.get("slides"), list):
        logger.warning("🚨 [GeneratePPT] No slides array found, returning None")
        return None

    if "meta" not in parsed or not isinstance(parsed.get("meta"), dict):
        parsed["meta"] = {
            "title": "Generated Presentation",
            "author": "AI",
            "date": datetime.now().strftime("%Y-%m-%d"),
        }

    if "design" not in parsed or not isinstance(parsed.get("design"), dict):
        parsed["design"] = {
            "palette": {
                "primary": "#1e40af",
                "secondary": "#3b82f6",
                "accent": "#f59e0b",
                "background": "#ffffff",
                "text": "#1f2937",
                "text_light": "#6b7280",
            },
            "fonts": {"heading": "Inter", "body": "Inter"},
            "type_scale": {"heading": 2.5, "body": 1.0, "caption": 0.75},
            "radius": "8px",
        }

    for i, slide in enumerate(parsed["slides"]):
        if not isinstance(slide, dict):
            continue
        if "layout" not in slide:
            slide["layout"] = "content"
        if "title" not in slide:
            slide["title"] = f"Slide {i + 1}"
        if "image_prompt" not in slide or not slide["image_prompt"]:
            slide["image_prompt"] = "none"
        if "image_search_keyword" not in slide or not slide["image_search_keyword"]:
            slide["image_search_keyword"] = ""

    logger.info(f"📋 [GeneratePPT] Forced raw dict with {len(parsed['slides'])} slides + defaults")
    return parsed


async def _inject_image_urls(presentation_data: Dict[str, Any]) -> None:
    slides = presentation_data.get("slides", [])
    if not slides:
        return

    FALLBACK_KEYWORDS = [
        "modern minimalist office workspace bright",
        "futuristic robot hand glowing network blue",
        "clean white architecture interior sunlight",
        "abstract geometric shapes soft lighting 3d",
        "glowing charts screen macro close-up",
        "mechanical arm holographic network dark",
        "professional desk computer coffee clean",
        "cyberpunk city skyline night neon lights",
    ]

    from app.core.config import settings
    unsplash_key = getattr(settings, "UNSPLASH_ACCESS_KEY", "")

    injected_count = 0
    for idx, slide in enumerate(slides):
        layout = slide.get("layout", "content")

        if layout == "closing":
            continue

        keyword = slide.get("image_search_keyword", "").strip()
        if not keyword:
            keyword = slide.get("image_prompt", "").strip()
        if not keyword or keyword.lower() in ["none", "null", ""] or len(keyword) < 3:
            keyword = FALLBACK_KEYWORDS[idx % len(FALLBACK_KEYWORDS)]

        slide["image_search_keyword"] = keyword

        image_url = None

        if unsplash_key:
            image_url = await _search_unsplash_image(keyword, unsplash_key)

        if not image_url:
            keyword_str = keyword.strip().replace(" ", ",")
            keywords_list = [k for k in keyword_str.split(",") if k][:4]
            short_keyword = ",".join(keywords_list)
            encoded_keyword = urllib.parse.quote(short_keyword)
            image_url = f"https://loremflickr.com/1024/768/{encoded_keyword}?random=1"
            logger.info(f"🖼️ [LoremFlickr] keyword='{short_keyword}' → loremflickr.com search")

        slide["image_url"] = image_url
        injected_count += 1

        if layout not in ["cover", "section"]:
            if "components" not in slide or not isinstance(slide["components"], list):
                slide["components"] = []
            has_image_comp = any(isinstance(c, dict) and c.get("type") == "image" for c in slide["components"])
            if not has_image_comp:
                slide["components"].insert(0, {
                    "type": "image",
                    "image_url": image_url,
                })

    if injected_count > 0:
        logger.info(f"🖼️ [GeneratePPT] 注入了 {injected_count} 张配图 (keyword='{keyword}' per-slide)")


async def _search_unsplash_image(keyword: str, access_key: str) -> Optional[str]:
    try:
        import httpx
        encoded = urllib.parse.quote(keyword)
        url = f"https://api.unsplash.com/photos/random?query={encoded}&orientation=landscape"
        headers = {"Authorization": f"Client-ID {access_key}"}
        async with httpx.AsyncClient(timeout=5.0) as client:
            resp = await client.get(url, headers=headers)
            if resp.status_code == 200:
                data = resp.json()
                img_url = data.get("urls", {}).get("regular") or data.get("urls", {}).get("raw")
                if img_url:
                    logger.info(f"🖼️ [Unsplash] keyword='{keyword}' → {img_url[:80]}")
                    return img_url
            else:
                logger.warning(f"🖼️ [Unsplash] keyword='{keyword}' status={resp.status_code}")
    except Exception as e:
        logger.warning(f"🖼️ [Unsplash] search failed for '{keyword}': {e}")
    return None


async def _extract_document_content(file_path: str, file_type: str) -> dict:
    """
    提取文档内容，返回 {"text": str, "images": [data_uri, ...]}
    """
    if not file_path or not file_type:
        return {"text": "", "images": []}

    try:
        if file_type == "docx":
            from app.services.docx_parser import DocxParser
            parser = DocxParser()
            result = await parser.extract_text_with_images(file_path)
            text = result.get("text", "")
            if not text:
                html = await parser.convert_docx_to_html(file_path)
                if html:
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(html, "html.parser")
                    text = soup.get_text(separator="\n", strip=True)[:250000]
            images = result.get("images", [])
            return {"text": text[:250000], "images": images}

        elif file_type == "xlsx":
            from app.services.excel_service import ExcelDataProcessor
            processor = ExcelDataProcessor()
            schema = await processor.extract_excel_schema(file_path)
            return {"text": json.dumps(schema, ensure_ascii=False, indent=2)[:250000], "images": []}

        elif file_type == "pptx":
            text = ""
            ppt_images = []
            try:
                from app.services.ppt_service import PPTAnalyzer
                analyzer = PPTAnalyzer()
                result = await analyzer.extract_text_and_structure(file_path)
                text = result.get("full_text", "")
                ppt_images = result.get("images", [])
                if text:
                    text = text[:250000]
            except Exception as e:
                logger.warning(f"GeneratePPT: PPTAnalyzer failed: {e}, using safe fallback")

            if not text:
                try:
                    import pptx
                    prs = pptx.Presentation(file_path)
                    text_chunks = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                text_chunks.append(shape.text)
                    text = "\n".join(text_chunks)[:250000]
                except Exception as e:
                    logger.error(f"GeneratePPT: Safe PPTX fallback failed: {e}")

            return {"text": text, "images": ppt_images}

    except Exception as e:
        logger.warning(f"GeneratePPT: document extraction failed: {e}")

    return {"text": "", "images": []}


def _extract_search_queries(instruction: str, content: str) -> list:
    queries = []
    if instruction:
        queries.append(instruction[:200])
    if content:
        first_lines = content[:500].split("\n")
        for line in first_lines[:3]:
            line = line.strip()
            if line and len(line) > 5:
                queries.append(line[:100])
    return queries[:5] if queries else ["general"]


def _build_fallback_presentation(instruction: str, content: str) -> Dict[str, Any]:
    from app.models.ppt_schema import (
        DesignFonts,
        DesignPalette,
        DesignSystem,
        DesignTypeScale,
        SlideComponent,
        SlideMeta,
        SlidePage,
    )

    title = instruction[:60] if instruction else "AI Generated Presentation"
    content_preview = content[:500] if content else "No content available"

    deck = Presentation(
        meta=SlideMeta(title=title),
        design=DesignSystem(
            palette=DesignPalette(bg="#f7f5f0", text="#1a1814", accent="#6d4cff"),
            fonts=DesignFonts(),
            type_scale=DesignTypeScale(hero=168, body=36),
            radius=12,
        ),
        slides=[
            SlidePage(
                layout="cover",
                title=title,
                subtitle="Powered by Document Copilot",
                eyebrow="AI 自动生成",
                image_prompt="futuristic AI technology abstract blue glow",
                image_search_keyword="futuristic robot hand network",
                image_visual_description="发光的机械手触碰全息网络图，深蓝色背景，科技感十足",
            ),
            SlidePage(
                layout="content",
                title="内容概述",
                components=[
                    SlideComponent(
                        type="text",
                        content=content_preview[:300],
                    ),
                ],
                image_prompt="document paper stack desk",
                image_search_keyword="document paper stack desk",
                image_visual_description="整齐叠放的文件和纸张在木质桌面上，柔和自然光",
            ),
            SlidePage(
                layout="closing",
                title="谢谢",
            ),
        ],
    )
    return deck.model_dump()
